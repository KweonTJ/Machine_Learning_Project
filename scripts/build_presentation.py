#!/usr/bin/env python3
"""Build a Miricanvas-style 10-slide deck with python-pptx."""

from __future__ import annotations

import json
import struct
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = PROJECT_ROOT / "reports" / "final_project_presentation.pptx"

FONT = "Noto Sans CJK KR"

BG = RGBColor(247, 251, 255)
INK = RGBColor(20, 25, 38)
MUTED = RGBColor(96, 110, 128)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(225, 232, 240)

BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(8, 188, 218)
MINT = RGBColor(16, 185, 129)
LIME = RGBColor(132, 204, 22)
ORANGE = RGBColor(249, 115, 22)
PINK = RGBColor(236, 72, 153)
PURPLE = RGBColor(124, 58, 237)
YELLOW = RGBColor(245, 190, 20)

PALE_BLUE = RGBColor(232, 240, 255)
PALE_CYAN = RGBColor(226, 249, 252)
PALE_MINT = RGBColor(230, 248, 242)
PALE_LIME = RGBColor(241, 249, 226)
PALE_ORANGE = RGBColor(255, 238, 226)
PALE_PINK = RGBColor(255, 235, 245)
PALE_PURPLE = RGBColor(242, 236, 255)
PALE_YELLOW = RGBColor(255, 248, 220)


def read_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def f3(value: float) -> str:
    return f"{value:.3f}"


def png_size(path: Path) -> tuple[int, int]:
    with path.open("rb") as file:
        header = file.read(24)
    if header[:8] != b"\x89PNG\r\n\x1a\n":
        raise ValueError(f"Only PNG figures are supported: {path}")
    return struct.unpack(">II", header[16:24])


def fit_image(path: Path, x: float, y: float, w: float, h: float) -> tuple[float, float, float, float]:
    image_w, image_h = png_size(path)
    scale = min(w / image_w, h / image_h)
    fitted_w = image_w * scale
    fitted_h = image_h * scale
    return x + (w - fitted_w) / 2, y + (h - fitted_h) / 2, fitted_w, fitted_h


def shape(slide, x: float, y: float, w: float, h: float, fill: RGBColor, line: RGBColor | None = None, rounded: bool = False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    item = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    item.fill.solid()
    item.fill.fore_color.rgb = fill
    if line is None:
        item.line.fill.background()
    else:
        item.line.color.rgb = line
        item.line.width = Pt(0.6)
    return item


def text(
    slide,
    value: str | list[str],
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 12,
    color: RGBColor = INK,
    bold: bool = False,
    fill: RGBColor | None = None,
    line: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    bullet: bool = False,
    margin: float = 0.04,
    rounded: bool = False,
):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(0.6)

    frame = box.text_frame
    frame.clear()
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)

    lines = [value] if isinstance(value, str) else value
    for index, line_text in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.space_after = Pt(2)
        if bullet:
            line_text = f"• {line_text}"
        run = paragraph.add_run()
        run.text = line_text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def image(slide, image_path: str, x: float, y: float, w: float, h: float):
    path = PROJECT_ROOT / image_path
    ix, iy, iw, ih = fit_image(path, x, y, w, h)
    return slide.shapes.add_picture(str(path), Inches(ix), Inches(iy), width=Inches(iw), height=Inches(ih))


def base(slide, accent: RGBColor = BLUE) -> None:
    shape(slide, 0, 0, 13.333, 7.5, BG)
    shape(slide, 0.0, 0.0, 13.333, 0.12, accent)
    shape(slide, 12.58, 0.55, 0.34, 0.34, accent, None, True)
    shape(slide, 12.18, 0.55, 0.22, 0.22, YELLOW, None, True)


def slide(prs: Presentation, title: str, kicker: str, accent: RGBColor = BLUE):
    page = prs.slides.add_slide(prs.slide_layouts[6])
    base(page, accent)
    text(page, kicker, 0.62, 0.34, 4.2, 0.24, 8.5, accent, True, margin=0.01)
    text(page, title, 0.62, 0.62, 9.2, 0.54, 25, INK, True, margin=0.01)
    return page


def pill(slide, label: str, value: str, x: float, y: float, color: RGBColor, pale: RGBColor, w: float = 1.55) -> None:
    text(slide, value, x, y, w, 0.36, 15.5, color, True, pale, None, PP_ALIGN.CENTER, margin=0.02, rounded=True)
    text(slide, label, x, y + 0.39, w, 0.19, 7.0, MUTED, True, align=PP_ALIGN.CENTER, margin=0.01)


def small_figure(slide, caption: str, path: str, x: float, y: float, w: float, h: float, color: RGBColor) -> None:
    shape(slide, x, y, w, h, WHITE, LINE, True)
    shape(slide, x, y, w, 0.08, color)
    text(slide, caption, x + 0.14, y + 0.13, w - 0.28, 0.20, 8.0, MUTED, True, margin=0.01)
    image(slide, path, x + 0.12, y + 0.40, w - 0.24, h - 0.52)


def metric_card(slide, label: str, value: str, x: float, y: float, color: RGBColor, pale: RGBColor) -> None:
    shape(slide, x, y, 2.58, 0.74, pale, None, True)
    text(slide, value, x + 0.18, y + 0.10, 1.0, 0.26, 17, color, True, margin=0.01)
    text(slide, label, x + 0.18, y + 0.42, 2.12, 0.18, 7.6, MUTED, True, margin=0.01)


def dataset_card(slide, title: str, rows: str, target: str, model: str, x: float, y: float, color: RGBColor, pale: RGBColor) -> None:
    shape(slide, x, y, 3.78, 1.12, WHITE, LINE, True)
    shape(slide, x + 0.18, y + 0.19, 0.34, 0.34, pale, None, True)
    shape(slide, x + 0.27, y + 0.28, 0.16, 0.16, color, None, True)
    text(slide, title, x + 0.65, y + 0.14, 2.95, 0.22, 10.0, INK, True, margin=0.01)
    text(slide, rows, x + 0.65, y + 0.43, 1.05, 0.18, 7.5, color, True, margin=0.01)
    text(slide, target, x + 1.70, y + 0.43, 1.85, 0.18, 7.5, MUTED, True, margin=0.01)
    text(slide, model, x + 0.65, y + 0.73, 2.85, 0.20, 7.6, INK, fill=RGBColor(246, 248, 252), rounded=True, margin=0.02)


def process_card(slide, no: str, title: str, desc: str, x: float, color: RGBColor, pale: RGBColor) -> None:
    shape(slide, x, 1.58, 2.70, 1.86, WHITE, LINE, True)
    text(slide, no, x + 0.20, 1.84, 0.54, 0.35, 14.0, color, True, pale, align=PP_ALIGN.CENTER, margin=0.01, rounded=True)
    text(slide, title, x + 0.88, 1.80, 1.50, 0.24, 12.0, INK, True, margin=0.01)
    text(slide, desc, x + 0.27, 2.48, 2.16, 0.42, 8.8, MUTED, True, align=PP_ALIGN.CENTER, margin=0.01)


def result_slide(
    prs: Presentation,
    title: str,
    kicker: str,
    accent: RGBColor,
    metrics: list[tuple[str, str, RGBColor, RGBColor]],
    insight: list[str],
    figures: list[tuple[str, str, RGBColor]],
) -> None:
    page = slide(prs, title, kicker, accent)
    shape(page, 0.62, 1.38, 2.65, 5.60, WHITE, LINE, True)
    text(page, "핵심 해석", 0.90, 1.72, 2.0, 0.28, 14.0, accent, True, margin=0.01)
    y = 2.20
    for label, value, color, pale in metrics:
        metric_card(page, label, value, 0.92, y, color, pale)
        y += 0.86
    shape(page, 0.92, 4.90, 2.08, 0.03, RGBColor(230, 236, 242))
    text(page, insight, 0.90, 5.16, 2.15, 1.16, 8.7, INK, bullet=True, margin=0.01)

    small_figure(page, figures[0][0], figures[0][1], 3.68, 1.38, 4.02, 2.82, figures[0][2])
    small_figure(page, figures[1][0], figures[1][1], 8.05, 1.38, 4.02, 2.82, figures[1][2])
    small_figure(page, figures[2][0], figures[2][1], 3.68, 4.48, 3.10, 2.18, figures[2][2])
    small_figure(page, figures[3][0], figures[3][1], 7.12, 4.48, 4.95, 2.18, figures[3][2])


def build_deck() -> Presentation:
    metrics = {
        "mendeley": read_json(PROJECT_ROOT / "outputs" / "mendeley_metrics.json"),
        "obd_can": read_json(PROJECT_ROOT / "outputs" / "obd_can_metrics.json"),
        "vehicle_maintenance": read_json(PROJECT_ROOT / "outputs" / "vehicle_maintenance_metrics.json"),
        "automotive_obd_ii": read_json(PROJECT_ROOT / "outputs" / "automotive_obd_ii_metrics.json"),
        "ai4i": read_json(PROJECT_ROOT / "outputs" / "ai4i_metrics.json"),
    }

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    page = prs.slides.add_slide(prs.slide_layouts[6])
    shape(page, 0, 0, 13.333, 7.5, BG)
    shape(page, 7.45, 0, 5.88, 7.5, PALE_BLUE)
    shape(page, 0.62, 0.58, 1.30, 0.13, ORANGE)
    shape(page, 2.06, 0.58, 1.30, 0.13, CYAN)
    shape(page, 3.50, 0.58, 1.30, 0.13, LIME)
    text(page, "AUTOMOTIVE AI PROJECT", 0.64, 1.08, 2.4, 0.24, 8.5, BLUE, True, PALE_BLUE, align=PP_ALIGN.CENTER, margin=0.02, rounded=True)
    text(page, "AI 차량 점검\n리포트 시스템", 0.62, 1.62, 4.90, 1.34, 33, INK, True, margin=0.01)
    text(page, "차량 센서 데이터 학습 -> 새 데이터 적용 -> 결과 도출", 0.66, 3.18, 4.65, 0.34, 12.0, MUTED, True, margin=0.01)
    text(page, ["전처리", "머신러닝 알고리즘", "성능 평가", "시각적 결과물"], 0.66, 3.82, 2.72, 1.04, 10.4, INK, bullet=True, margin=0.01)
    pill(page, "OBD/CAN F1", f3(metrics["obd_can"]["macro_f1"]), 0.66, 5.62, MINT, PALE_MINT, 1.40)
    pill(page, "Mendeley F1", f3(metrics["mendeley"]["macro_f1"]), 2.28, 5.62, CYAN, PALE_CYAN, 1.40)
    pill(page, "정비 F1", f3(metrics["vehicle_maintenance"]["macro_f1"]), 3.90, 5.62, ORANGE, PALE_ORANGE, 1.40)

    shape(page, 6.28, 0.70, 6.45, 5.95, WHITE, LINE, True)
    small_figure(page, "OBD/CAN Confusion", "outputs/figures/obd_can_confusion_matrix.png", 6.68, 1.05, 2.75, 2.35, MINT)
    small_figure(page, "Feature Importance", "outputs/figures/mendeley_feature_importance.png", 9.72, 1.05, 2.28, 2.35, BLUE)
    small_figure(page, "Maintenance Distribution", "outputs/figures/vehicle_maintenance_feature_distribution.png", 6.68, 3.70, 3.55, 2.20, ORANGE)
    small_figure(page, "AI4I Confusion", "outputs/figures/ai4i_confusion_matrix.png", 10.54, 3.70, 1.46, 2.20, PURPLE)

    page = slide(prs, "프로젝트 핵심 흐름", "DATA TO DECISION", BLUE)
    for no, title, desc, x, color, pale in [
        ("01", "Raw data", "CSV/ZIP 센서·정비 데이터 수집", 0.70, CYAN, PALE_CYAN),
        ("02", "Preprocess", "컬럼 표준화, 결측 제거, 라벨 매핑", 3.62, BLUE, PALE_BLUE),
        ("03", "Feature", "센서 magnitude, trip 통계, one-hot", 6.54, YELLOW, PALE_YELLOW),
        ("04", "ML & Report", "분류·이상 탐지 후 그래프 생성", 9.46, MINT, PALE_MINT),
    ]:
        process_card(page, no, title, desc, x, color, pale)
    for x, color in [(3.22, BLUE), (6.14, YELLOW), (9.06, MINT)]:
        text(page, "→", x, 2.34, 0.26, 0.24, 18, color, True, align=PP_ALIGN.CENTER, margin=0.01)
    shape(page, 0.85, 4.25, 5.62, 1.58, PALE_BLUE, None, True)
    text(page, "프로젝트 포인트", 1.16, 4.55, 5.0, 0.28, 15, BLUE, True, margin=0.01)
    text(page, "학습 데이터만 맞추는 것이 아니라, 새 데이터에도 같은 전처리와 feature schema를 적용해 결과를 낼 수 있게 구성했습니다.", 1.16, 5.04, 4.92, 0.40, 10.0, INK, margin=0.01)
    shape(page, 6.90, 4.25, 5.54, 1.58, PALE_ORANGE, None, True)
    text(page, "산출물", 7.20, 4.55, 4.8, 0.28, 15, ORANGE, True, margin=0.01)
    text(page, "metrics JSON, confusion matrix, feature importance, F1, distribution 그래프와 최종 PDF/PPT를 생성했습니다.", 7.20, 5.04, 4.78, 0.40, 10.0, INK, margin=0.01)

    page = slide(prs, "사용 데이터셋", "FIVE EXPERIMENTS + ONE VERIFIED DATASET", CYAN)
    dataset_card(page, "Mendeley phone sensor", "14,249 rows", "normal/aggressive", "standardized k-NN", 0.68, 1.36, CYAN, PALE_CYAN)
    dataset_card(page, "OBD-II/CAN Driving", "555,000 rows", "moderate/aggressive", "standardized class centroid", 4.75, 1.36, MINT, PALE_MINT)
    dataset_card(page, "Vehicle Maintenance", "1,970 rows", "normal/issue", "robust z-score anomaly", 8.82, 1.36, ORANGE, PALE_ORANGE)
    dataset_card(page, "KIT Automotive OBD-II", "81 trip CSV", "normal/free_flow/traffic", "trip-level class centroid", 0.68, 3.05, BLUE, PALE_BLUE)
    dataset_card(page, "AI4I 2020", "10,000 rows", "normal/failure", "predictive maintenance benchmark", 4.75, 3.05, PURPLE, PALE_PURPLE)
    dataset_card(page, "APS Failure", "76,000 rows", "download verified", "excluded from final model", 8.82, 3.05, PINK, PALE_PINK)
    shape(page, 0.78, 5.52, 11.72, 0.84, PALE_YELLOW, None, True)
    text(page, "APS는 다운로드와 무결성은 확인했지만 feature가 익명화되어 차량 상태 설명력이 낮아 최종 모델에서는 제외했습니다.", 1.08, 5.76, 11.0, 0.23, 10.5, INK, True, margin=0.01)

    page = slide(prs, "전처리와 알고리즘 활용", "WHAT HAPPENS BEFORE TRAINING", ORANGE)
    shape(page, 0.75, 1.34, 5.70, 5.00, WHITE, LINE, True)
    text(page, "학습 전 데이터 처리", 1.10, 1.74, 4.8, 0.34, 16, ORANGE, True, margin=0.01)
    text(
        page,
        [
            "원천 CSV/ZIP 로드 후 표준 feature명으로 컬럼 매핑",
            "수치형 변환, 필수 feature와 라벨 기준 결측 행 제거",
            "OBD Label 0->aggressive, Label 1->moderate",
            "정비 데이터 failure flag를 issue_label로 통합",
            "train/test 75/25 stratified split, seed=42",
        ],
        1.10,
        2.32,
        4.75,
        2.05,
        9.8,
        INK,
        bullet=True,
        margin=0.01,
    )
    shape(page, 6.78, 1.34, 5.70, 5.00, WHITE, LINE, True)
    text(page, "머신러닝 활용", 7.13, 1.74, 4.8, 0.34, 16, BLUE, True, margin=0.01)
    for label, body, y, color, pale in [
        ("k-NN", "Mendeley 센서를 표준화하고 가까운 학습 샘플 다수결로 주행 습관 분류", 2.32, CYAN, PALE_CYAN),
        ("Class centroid", "클래스별 평균 벡터를 학습하고 새 데이터는 가장 가까운 중심으로 분류", 3.50, BLUE, PALE_BLUE),
        ("Robust z-score", "정상 정비 데이터의 median/MAD로 정상 범위 이탈을 issue로 판정", 4.68, ORANGE, PALE_ORANGE),
    ]:
        text(page, label, 7.13, y, 1.50, 0.30, 9.8, color, True, pale, align=PP_ALIGN.CENTER, margin=0.02, rounded=True)
        text(page, body, 8.86, y - 0.02, 3.02, 0.34, 8.4, INK, margin=0.01)

    result_slide(
        prs,
        "Mendeley 스마트폰 센서",
        "RESULT 01",
        CYAN,
        [
            ("Accuracy", f3(metrics["mendeley"]["accuracy"]), CYAN, PALE_CYAN),
            ("Macro F1", f3(metrics["mendeley"]["macro_f1"]), MINT, PALE_MINT),
            ("Train rows", str(metrics["mendeley"]["train_rows"]), BLUE, PALE_BLUE),
        ],
        ["가속도/자이로 기반 분류", "standardized k-NN 적용", "센서 기반 기준선 확보"],
        [
            ("Confusion Matrix", "outputs/figures/mendeley_confusion_matrix.png", CYAN),
            ("Feature Importance", "outputs/figures/mendeley_feature_importance.png", BLUE),
            ("F1 Scores", "outputs/figures/mendeley_f1_scores.png", ORANGE),
            ("Sensor Distribution", "outputs/figures/mendeley_sensor_distribution.png", MINT),
        ],
    )

    result_slide(
        prs,
        "OBD-II/CAN 운전 습관",
        "RESULT 02",
        MINT,
        [
            ("Accuracy", f3(metrics["obd_can"]["accuracy"]), CYAN, PALE_CYAN),
            ("Macro F1", f3(metrics["obd_can"]["macro_f1"]), MINT, PALE_MINT),
            ("Train rows", str(metrics["obd_can"]["train_rows"]), BLUE, PALE_BLUE),
        ],
        ["차량 내부 신호 사용", "class centroid 분류", "가장 안정적인 결과"],
        [
            ("Confusion Matrix", "outputs/figures/obd_can_confusion_matrix.png", MINT),
            ("Feature Importance", "outputs/figures/obd_can_feature_importance.png", BLUE),
            ("F1 Scores", "outputs/figures/obd_can_f1_scores.png", ORANGE),
            ("Feature Distribution", "outputs/figures/obd_can_feature_distribution.png", YELLOW),
        ],
    )

    result_slide(
        prs,
        "Vehicle Maintenance 이상 탐지",
        "RESULT 03",
        ORANGE,
        [
            ("Accuracy", f3(metrics["vehicle_maintenance"]["accuracy"]), CYAN, PALE_CYAN),
            ("Macro F1", f3(metrics["vehicle_maintenance"]["macro_f1"]), ORANGE, PALE_ORANGE),
            ("Threshold", f3(metrics["vehicle_maintenance"]["threshold"]), BLUE, PALE_BLUE),
        ],
        ["정상 telemetry 기준 학습", "정상 범위 이탈을 issue로 판정", "issue recall 1.000"],
        [
            ("Confusion Matrix", "outputs/figures/vehicle_maintenance_confusion_matrix.png", ORANGE),
            ("Feature Importance", "outputs/figures/vehicle_maintenance_feature_importance.png", BLUE),
            ("F1 Scores", "outputs/figures/vehicle_maintenance_f1_scores.png", MINT),
            ("Feature Distribution", "outputs/figures/vehicle_maintenance_feature_distribution.png", YELLOW),
        ],
    )

    result_slide(
        prs,
        "KIT Automotive OBD-II",
        "RESULT 04",
        BLUE,
        [
            ("Accuracy", f3(metrics["automotive_obd_ii"]["accuracy"]), CYAN, PALE_CYAN),
            ("Macro F1", f3(metrics["automotive_obd_ii"]["macro_f1"]), PINK, PALE_PINK),
            ("Trip files", str(metrics["automotive_obd_ii"]["train_rows"] + metrics["automotive_obd_ii"]["test_rows"]), BLUE, PALE_BLUE),
        ],
        ["OBD 로그를 trip 단위 집계", "도로 조건 3-class 분류", "추가 데이터 필요"],
        [
            ("Confusion Matrix", "outputs/figures/automotive_obd_ii_confusion_matrix.png", BLUE),
            ("Feature Importance", "outputs/figures/automotive_obd_ii_feature_importance.png", CYAN),
            ("F1 Scores", "outputs/figures/automotive_obd_ii_f1_scores.png", ORANGE),
            ("Feature Distribution", "outputs/figures/automotive_obd_ii_feature_distribution.png", YELLOW),
        ],
    )

    result_slide(
        prs,
        "AI4I 예지정비 벤치마크",
        "RESULT 05",
        PURPLE,
        [
            ("Accuracy", f3(metrics["ai4i"]["accuracy"]), CYAN, PALE_CYAN),
            ("Macro F1", f3(metrics["ai4i"]["macro_f1"]), ORANGE, PALE_ORANGE),
            ("Train rows", str(metrics["ai4i"]["train_rows"]), PURPLE, PALE_PURPLE),
        ],
        ["온도·토크·마모 feature", "예지정비 구조 검증", "failure class 개선 필요"],
        [
            ("Confusion Matrix", "outputs/figures/ai4i_confusion_matrix.png", PURPLE),
            ("Feature Importance", "outputs/figures/ai4i_feature_importance.png", BLUE),
            ("F1 Scores", "outputs/figures/ai4i_f1_scores.png", ORANGE),
            ("Feature Distribution", "outputs/figures/ai4i_feature_distribution.png", MINT),
        ],
    )

    page = slide(prs, "최종 결론", "SCORECARD", MINT)
    shape(page, 0.75, 1.36, 5.56, 4.58, WHITE, LINE, True)
    text(page, "Macro F1 비교", 1.08, 1.73, 4.8, 0.32, 16, INK, True, margin=0.01)
    score_rows = [
        ("OBD-II/CAN", metrics["obd_can"]["macro_f1"], MINT),
        ("Mendeley", metrics["mendeley"]["macro_f1"], CYAN),
        ("Vehicle Maintenance", metrics["vehicle_maintenance"]["macro_f1"], ORANGE),
        ("AI4I", metrics["ai4i"]["macro_f1"], PURPLE),
        ("KIT OBD-II", metrics["automotive_obd_ii"]["macro_f1"], BLUE),
    ]
    y = 2.36
    for label, value, color in score_rows:
        text(page, label, 1.08, y - 0.02, 1.75, 0.22, 8.7, INK, True, margin=0.01)
        shape(page, 2.95, y + 0.04, 2.32, 0.18, RGBColor(230, 236, 242), None, True)
        shape(page, 2.95, y + 0.04, 2.32 * value, 0.18, color, None, True)
        text(page, f3(value), 5.40, y - 0.03, 0.54, 0.22, 8.5, color, True, margin=0.01)
        y += 0.57
    shape(page, 6.75, 1.36, 5.72, 1.72, PALE_MINT, None, True)
    text(page, "핵심 성과", 7.08, 1.70, 4.8, 0.30, 16, MINT, True, margin=0.01)
    text(page, ["여러 데이터셋을 같은 평가 흐름으로 비교", "분류와 이상 탐지를 함께 활용", "그래프와 보고서 자동 생성"], 7.08, 2.17, 4.65, 0.70, 9.5, INK, bullet=True, margin=0.01)
    shape(page, 6.75, 3.42, 5.72, 1.72, PALE_ORANGE, None, True)
    text(page, "다음 단계", 7.08, 3.76, 4.8, 0.30, 16, ORANGE, True, margin=0.01)
    text(page, ["KIT OBD-II 추가 라벨 데이터 확보", "실제 차량 고장/정비 라벨 확장", "시계열 모델과 실시간 적용 구조로 개선"], 7.08, 4.23, 4.65, 0.70, 9.5, INK, bullet=True, margin=0.01)
    shape(page, 0.78, 6.35, 11.70, 0.55, PALE_BLUE, None, True)
    text(page, "수업자료 연결: 데이터 전처리, 지도학습 분류, 이상 탐지, 성능 평가, 시각화를 실제 프로젝트 흐름으로 구현", 1.08, 6.50, 10.95, 0.20, 10.2, BLUE, True, margin=0.01)

    return prs


def main() -> int:
    prs = build_deck()
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"created: {OUTPUT_PATH}")
    print(f"slides: {len(prs.slides)}")
    print("unique source figures embedded: 20")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
